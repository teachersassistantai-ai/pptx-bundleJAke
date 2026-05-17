"""pptx_bundle -- generated; do not edit. Bundles pptx_builder for Relevance.ai."""

# === brand.py ===
"""resolve_brand -- turn generative design_tokens into an internal brand dict.

The renderer is fully token-driven. `design_tokens` is the dict described by
DESIGN_TOKEN_SCHEMA in src/rai_builder/schema.py (palette / type / composition
/ motif). `resolve_brand` converts those categorical token values into the
concrete numbers (colours, point sizes, inch margins, gaps) the layout
renderers consume. Nothing here is hardcoded per deck -- two different
design_tokens inputs produce visibly different brand dicts.

This module is a plain-function + dict export so it survives clean
concatenation into the Relevance.ai bundle.
"""

SLIDE_W = 13.333
SLIDE_H = 7.5

# Fallback hexes used only when a palette role is missing/malformed.
_FALLBACK_PALETTE = {
    "bg": "FFFFFF",
    "ink": "1A2747",
    "accent": "2EC4B6",
    "accent_2": "E76F51",
    "muted": "55617A",
    "surface": "F1F1F1",
}

# type.scale -> a base point-size set. Sizes step up across the three scales.
_TYPE_SCALES = {
    "small": {
        "title": 40, "heading": 24, "statement": 32, "body": 14,
        "caption": 9, "stat": 64, "pillar_heading": 16, "section_number": 72,
    },
    "medium": {
        "title": 54, "heading": 32, "statement": 44, "body": 18,
        "caption": 11, "stat": 88, "pillar_heading": 20, "section_number": 96,
    },
    "large": {
        "title": 66, "heading": 40, "statement": 56, "body": 22,
        "caption": 13, "stat": 108, "pillar_heading": 24, "section_number": 120,
    },
}

# composition.margin -> outer inch margin.
_MARGINS = {"narrow": 0.5, "standard": 0.75, "wide": 1.0}

# composition.density -> body-size delta + content breathing room.
_DENSITY = {
    "sparse":   {"body_delta": 2,  "block_gap": 0.35},
    "balanced": {"body_delta": 0,  "block_gap": 0.22},
    "dense":    {"body_delta": -2, "block_gap": 0.12},
}

# composition.whitespace -> inter-block gap multiplier.
_WHITESPACE = {"tight": 0.7, "generous": 1.3}

# type.tracking -> approximate spacing nudge (python-pptx has no real tracking).
_TRACKING = {"tight": -0.3, "normal": 0.0, "wide": 0.6}

# motif.weight -> relative thickness/scale of the motif element.
_MOTIF_WEIGHT = {"light": 1.0, "bold": 2.2}


def _pick(d, value, default_key):
    """Safe categorical lookup -- unknown value falls back, never raises."""
    if value in d:
        return d[value]
    return d[default_key]


def _clean_hex(value, fallback):
    """Coerce a hex-ish string to a 6-char hex with no leading #."""
    if not isinstance(value, str):
        return fallback
    h = value.strip().lstrip("#")
    if len(h) == 3:  # shorthand -> expand
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return fallback
    try:
        int(h, 16)
    except ValueError:
        return fallback
    return h.upper()


def resolve_brand(design_tokens):
    """Turn design_tokens into the internal brand dict layout renderers consume.

    Args:
        design_tokens: dict matching DESIGN_TOKEN_SCHEMA. Missing keys are
            tolerated -- every lookup has a defensive default.
    Returns:
        brand dict with keys: colours, fonts, header_case, tracking, sizes_pt,
        composition, motif, slide_dimensions_in.
    """
    tokens = design_tokens if isinstance(design_tokens, dict) else {}
    palette = tokens.get("palette", {}) or {}
    type_t = tokens.get("type", {}) or {}
    comp = tokens.get("composition", {}) or {}
    motif = tokens.get("motif", {}) or {}

    colours = {
        role: _clean_hex(palette.get(role), _FALLBACK_PALETTE[role])
        for role in _FALLBACK_PALETTE
    }

    scale_key = type_t.get("scale", "medium")
    sizes = dict(_pick(_TYPE_SCALES, scale_key, "medium"))
    density_key = comp.get("density", "balanced")
    density = _pick(_DENSITY, density_key, "balanced")
    # Density nudges body / pillar text without touching display sizes.
    sizes["body"] = max(9, sizes["body"] + density["body_delta"])
    sizes["pillar_heading"] = max(11, sizes["pillar_heading"] + density["body_delta"])

    margin = _pick(_MARGINS, comp.get("margin", "standard"), "standard")
    whitespace_mult = _pick(_WHITESPACE, comp.get("whitespace", "generous"), "generous")
    block_gap = round(density["block_gap"] * whitespace_mult, 3)

    return {
        "colours": colours,
        "fonts": {
            "header": type_t.get("header_font", "Georgia"),
            "body": type_t.get("body_font", "Calibri"),
        },
        "header_case": type_t.get("header_case", "title"),
        "tracking": _pick(_TRACKING, type_t.get("tracking", "normal"), "normal"),
        "sizes_pt": sizes,
        "composition": {
            "margin": margin,
            "alignment": comp.get("alignment", "left"),
            "density": density_key,
            "block_gap": block_gap,
            "whitespace": comp.get("whitespace", "generous"),
        },
        "motif": {
            "kind": motif.get("kind", "none"),
            "weight": _pick(_MOTIF_WEIGHT, motif.get("weight", "light"), "light"),
        },
        "slide_dimensions_in": {"width": SLIDE_W, "height": SLIDE_H},
    }

# === palette.py ===
"""Token-driven geometry.

The old renderer baked every box coordinate into a static PALETTE dict. This
module replaces that: `compute_geometry(brand)` derives all box geometry from
`brand["composition"]` (margins, block gap, alignment). Two different
design_tokens -> two different geometries (margins shift, blocks move).

All measurements in inches. Slide is 13.333 x 7.5 (16:9).
"""



def compute_geometry(brand):
    """Return a geometry dict the layout renderers consume.

    Everything is a function of the resolved composition tokens. The renderer
    never reads hardcoded coordinates -- it reads this.
    """
    comp = brand["composition"]
    m = comp["margin"]                       # outer margin (in)
    gap = comp["block_gap"]                  # inter-block gap (in)
    align = comp["alignment"]                # left | centered | asymmetric

    content_w = SLIDE_W - 2 * m
    content_h = SLIDE_H - 2 * m

    # Asymmetric layouts nudge the content origin right, creating a wider
    # left margin -- a visible compositional signature.
    asym = 0.0
    if align == "asymmetric":
        asym = round(m * 0.9, 3)

    body_left = m + asym
    body_w = content_w - asym

    # Standard heading band sits at the top, content begins below it.
    heading_h = 0.95
    heading_top = m
    content_top = heading_top + heading_h + gap

    return {
        "slide_w": SLIDE_W,
        "slide_h": SLIDE_H,
        "margin": m,
        "gap": gap,
        "alignment": align,
        "content_w": content_w,
        "content_h": content_h,
        "body_left": body_left,
        "body_w": body_w,
        "heading_top": heading_top,
        "heading_h": heading_h,
        "content_top": content_top,
        "content_bottom": SLIDE_H - m,
        "asym": asym,
    }


def heading_box(geom):
    """The standard top heading band (left, top, width, height)."""
    return (geom["margin"], geom["heading_top"], geom["content_w"], geom["heading_h"])


def body_region(geom):
    """The region below the heading band: (left, top, width, height)."""
    top = geom["content_top"]
    return (geom["body_left"], top, geom["body_w"], geom["content_bottom"] - top)

# === helpers.py ===
"""Shared utilities for the token-driven layout renderers.

Everything visual flows through here: colours come from `brand["colours"]`,
fonts from `brand["fonts"]`, point sizes from `brand["sizes_pt"]`. Layout
renderers never hardcode a hex, a family, or a size.

Also home to:
  * header-case transform (type.header_case)
  * count-guard hardening (VARIANT_ELEMENT_COUNTS, iteration-0 carry-forward)
  * motif drawing (all five motif kinds)
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# --- colour -----------------------------------------------------------------

def hex_to_rgb(hex_str):
    h = (hex_str or "000000").lstrip("#")
    if len(h) != 6:
        h = "000000"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def colour(brand, role):
    """Resolve a palette role to a hex string. Unknown role -> ink."""
    return brand["colours"].get(role, brand["colours"]["ink"])


# Semantic accent_role -> palette role mapping. Used everywhere a slide spec
# carries an `accent_role` field. neutral and positive use the dominant accent;
# warning and critical escalate to accent_2. This is the single source of truth
# for "what does accent_role mean visually" -- every layout reads through this.
_ROLE_TO_PALETTE = {
    "neutral": "accent",
    "positive": "accent",
    "warning": "accent_2",
    "critical": "accent_2",
}


def accent_for_role(brand, role):
    """Resolve a slide-level accent_role to a hex string from the brand palette.

    Unknown role -> the dominant accent. This keeps the visual deck-coherent
    while letting individual slides escalate to accent_2 for warnings/criticals.
    """
    palette_role = _ROLE_TO_PALETTE.get(role, "accent")
    return colour(brand, palette_role)


def _relative_luminance(hex_str):
    """WCAG-ish relative luminance for contrast-aware text colour choice."""
    h = (hex_str or "000000").lstrip("#")
    if len(h) != 6:
        return 0.0
    chan = []
    for i in (0, 2, 4):
        c = int(h[i:i + 2], 16) / 255.0
        c = c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
        chan.append(c)
    return 0.2126 * chan[0] + 0.7152 * chan[1] + 0.0722 * chan[2]


def readable_on(brand, fill_hex):
    """Pick bg or ink (whichever contrasts the fill better) for text on a fill."""
    fill_l = _relative_luminance(fill_hex)
    ink_l = _relative_luminance(colour(brand, "ink"))
    bg_l = _relative_luminance(colour(brand, "bg"))

    def ratio(a, b):
        hi, lo = max(a, b), min(a, b)
        return (hi + 0.05) / (lo + 0.05)

    return colour(brand, "ink") if ratio(fill_l, ink_l) >= ratio(fill_l, bg_l) \
        else colour(brand, "bg")


# --- text -------------------------------------------------------------------

def apply_header_case(brand, text):
    """Apply type.header_case to header text at render time."""
    text = "" if text is None else str(text)
    case = brand.get("header_case", "title")
    if case == "upper":
        return text.upper()
    if case == "sentence":
        return text[:1].upper() + text[1:] if text else text
    if case == "title":
        # Title-case but keep short connectors lowercase for a designed look.
        small = {"a", "an", "the", "of", "and", "or", "to", "in", "on", "at",
                 "for", "with", "by", "vs", "via", "per", "from", "into", "but"}
        words = text.split()
        out = []
        for idx, w in enumerate(words):
            lw = w.lower()
            if idx != 0 and lw in small:
                out.append(lw)
            else:
                out.append(w[:1].upper() + w[1:].lower() if w else w)
        return " ".join(out)
    return text


def _fit_font_size(text, width_in, height_in, size_pt, line_spacing=1.15):
    """Estimate the largest font size (<= size_pt) at which `text` fits a box
    of width_in x height_in inches.

    python-pptx text boxes do not auto-shrink, so long LLM-written content
    overflows fixed-height boxes (callout bodies, quadrant cells, headings,
    long titles). This heuristic -- proportional-font average glyph advance
    ~0.52em, plus a 1.2 leading fudge -- picks a size that prevents gross
    overflow. It is deliberately approximate; its job is "never spill off the
    box", not pixel-perfect typesetting. Floors at 55% of the requested size
    (or 9pt) so display type still reads big.
    """
    text = "" if text is None else str(text)
    if not text.strip() or width_in <= 0 or height_in <= 0:
        return size_pt
    width_pt = width_in * 72.0
    height_pt = height_in * 72.0
    paras = text.split("\n")
    floor = max(9.0, size_pt * 0.55)
    size = float(size_pt)
    while size > floor:
        chars_per_line = max(1, int(width_pt / (0.52 * size)))
        lines = 0
        for para in paras:
            n = len(para)
            lines += 1 if n == 0 else -(-n // chars_per_line)
        needed = lines * size * line_spacing * 1.2
        if needed <= height_pt:
            return round(size, 1)
        size -= 1.0
    return round(floor, 1)


def add_text_box(
    slide, brand,
    left_in, top_in, width_in, height_in,
    text,
    font_role,            # "header" | "body"
    size_pt,
    colour_hex,
    align="left",
    anchor="top",
    bold=False,
    is_header=False,
):
    """Add a single-paragraph text box, fully token-driven.

    `font_role` selects the family from brand["fonts"]. `is_header` triggers
    the header-case transform.
    """
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.TOP)

    p = tf.paragraphs[0]
    rendered = apply_header_case(brand, text) if is_header else ("" if text is None else str(text))
    p.text = rendered
    p.alignment = {
        "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    family = brand["fonts"].get(font_role, brand["fonts"]["body"])
    fitted = _fit_font_size(rendered, width_in, height_in, size_pt)
    if p.runs:
        run = p.runs[0]
        run.font.name = family
        run.font.size = Pt(fitted)
        run.font.color.rgb = hex_to_rgb(colour_hex)
        run.font.bold = bold
        # Approximate tracking via character spacing where pptx allows it.
        try:
            run.font._rPr.set("spc", str(int(Pt(brand.get("tracking", 0.0)))))
        except Exception:
            pass
    return tb


def add_multiline(
    slide, brand,
    left_in, top_in, width_in, height_in,
    lines,
    font_role, size_pt, colour_hex,
    align="left", anchor="top", bullet=False,
    bullet_colour_hex=None, line_spacing=1.15,
):
    """Add a multi-paragraph text box (one paragraph per line).

    bullet=True draws a real round-bullet glyph (U+2022) via OOXML paragraph
    properties, not "- text". bullet_colour_hex (optional) lets the bullet sit
    in the accent colour while the line text stays in colour_hex -- a small
    but high-impact design lift on list slides.

    line_spacing controls inter-paragraph leading. Default 1.15 = comfortable
    body-text spacing; sparse decks benefit from 1.3+.
    """
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.TOP)
    family = brand["fonts"].get(font_role, brand["fonts"]["body"])
    pp_align = {
        "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    items = list(lines) if lines else [""]
    # Fit one size across all paragraphs so the list reads as one block.
    _joined = "\n".join("" if x is None else str(x) for x in items)
    fitted = _fit_font_size(_joined, width_in, height_in, size_pt, line_spacing)
    for idx, line in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = "" if line is None else str(line)
        p.alignment = pp_align
        if p.runs:
            run = p.runs[0]
            run.font.name = family
            run.font.size = Pt(fitted)
            run.font.color.rgb = hex_to_rgb(colour_hex)
        # Real bullet via OOXML: <a:pPr><a:buClr/><a:buChar char="•"/></a:pPr>.
        # The pPr element exists by default; we just add buChar (and optional
        # buClr) children. Falls back silently if the XML hooks aren't there.
        if bullet:
            try:
                from pptx.oxml.ns import qn
                pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
                # Indent so the bullet sits in negative space.
                pPr.set("marL", str(int(228600)))  # ~0.25 inch
                pPr.set("indent", str(int(-228600)))
                # Bullet colour (optional).
                if bullet_colour_hex:
                    buClr = pPr.makeelement(qn("a:buClr"), {})
                    srgb = pPr.makeelement(qn("a:srgbClr"),
                                           {"val": bullet_colour_hex.lstrip("#")})
                    buClr.append(srgb)
                    pPr.append(buClr)
                # Bullet glyph.
                buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
                pPr.append(buChar)
            except Exception:
                # If the XML manipulation fails (older pptx), fall back to
                # a Unicode bullet inline rather than the old ASCII hyphen.
                p.text = "\u2022  " + (p.text or "")
        # Line spacing -- paragraph-level setting that python-pptx exposes.
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
    return tb


# --- shapes -----------------------------------------------------------------

def _add_shape(slide, shape_enum, left, top, width, height, fill_hex, line_hex=None):
    shape = slide.shapes.add_shape(
        shape_enum, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
        shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    return shape


def add_filled_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    return _add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height,
                      fill_hex, line_hex)


def add_rounded_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    return _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width,
                      height, fill_hex, line_hex)


def add_oval(slide, left, top, width, height, fill_hex, line_hex=None):
    return _add_shape(slide, MSO_SHAPE.OVAL, left, top, width, height,
                      fill_hex, line_hex)


def fill_background(slide, brand):
    """Paint the slide background with the palette bg colour."""
    geom = brand["slide_dimensions_in"]
    rect = add_filled_rect(slide, 0, 0, geom["width"], geom["height"],
                           colour(brand, "bg"))
    # Send to back so subsequent shapes layer on top.
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return rect


# --- heading band (shared across content layouts) ---------------------------

def draw_heading(slide, brand, geom, text, accent_role="neutral",
                 with_rule=True):
    """Draw a content slide's heading: title text + an accent rule beneath it.

    Replaces the hand-rolled "add_text_box for heading" call that every content
    layout was duplicating. Centralising it lets us add a designed accent rule
    consistently and pick the rule colour from accent_role.

    The rule is a thin filled rect just under the heading, accent-coloured,
    ~1.4 inches long. It gives every content slide a recognisable signature
    without dominating the composition.
    """
    sizes = brand["sizes_pt"]
    add_text_box(
        slide, brand, geom["margin"], geom["heading_top"],
        geom["content_w"], geom["heading_h"],
        as_text(text), "header", sizes["heading"], colour(brand, "ink"),
        bold=True, is_header=True,
    )
    if with_rule:
        rule_y = geom["heading_top"] + geom["heading_h"] + 0.02
        rule_colour = accent_for_role(brand, accent_role)
        add_filled_rect(slide, geom["margin"], rule_y, 1.4, 0.05, rule_colour)


# --- motif ------------------------------------------------------------------

def draw_motif(slide, brand):
    """Draw the per-deck motif on a content slide in the accent colour.

    Implements all five motif kinds. `motif.weight` scales thickness/size.
    A `none` kind is a safe no-op.
    """
    motif = brand.get("motif", {})
    kind = motif.get("kind", "none")
    if kind == "none":
        return
    weight = motif.get("weight", 1.0)
    accent = colour(brand, "accent")
    dims = brand["slide_dimensions_in"]
    w, h = dims["width"], dims["height"]

    if kind == "corner_block":
        size = round(0.75 * weight, 3)
        add_filled_rect(slide, w - size, 0, size, size, accent)

    elif kind == "side_rule":
        # Side rule is the most-used motif. Make it actually visible:
        # 0.18" at light weight, 0.40" at bold.
        bar = round(0.18 * weight, 3)
        add_filled_rect(slide, 0, 0, bar, h, accent)

    elif kind == "dot_grid":
        dot = round(0.10 * weight, 3)
        step = 0.42
        y = h - 1.0
        x = w - 1.6
        for row in range(3):
            for col in range(5):
                add_oval(slide, x + col * step, y + row * step, dot, dot, accent)

    elif kind == "baseline_rule":
        bar = round(0.08 * weight, 3)
        margin = brand["composition"]["margin"]
        add_filled_rect(slide, margin, h - margin + 0.1,
                        w - 2 * margin, bar, accent)


# --- count-guard hardening (iteration-0 carry-forward) ----------------------

def coerce_count(items, min_n, max_n, filler):
    """Pad/truncate a list to [min_n, max_n] -- never raises.

    `filler` is a zero-arg callable producing a placeholder element used to
    pad a short list.
    """
    out = list(items) if isinstance(items, (list, tuple)) else []
    if len(out) > max_n:
        out = out[:max_n]
    while len(out) < min_n:
        out.append(filler())
    return out


def coerce_int(value, default=0):
    """Coerce a possibly-string value to int -- safe for `:02d` formatting."""
    if isinstance(value, bool):
        return default
    if isinstance(value, int):
        return value
    if isinstance(value, float):
        return int(value)
    if isinstance(value, str):
        digits = value.strip().lstrip("+-")
        if digits.isdigit():
            n = int(digits)
            return -n if value.strip().startswith("-") else n
    return default


def as_text(value, default=""):
    """Coerce any value to a string for safe rendering."""
    if value is None:
        return default
    if isinstance(value, str):
        return value
    if isinstance(value, (int, float, bool)):
        return str(value)
    return default

# === layouts/title.py ===
"""statement family / title variant: typography hero + accent edge block."""



def _render_title(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    align = "center" if geom["alignment"] == "centered" else "left"

    # Accent edge block -- left rail, full height. Width scales with motif.
    rail = round(0.35 + 0.15 * brand["motif"]["weight"], 3)
    add_filled_rect(slide, 0, 0, rail, geom["slide_h"], colour(brand, "accent"))

    left = geom["body_left"]
    width = geom["body_w"]

    # Short accent rule above the title -- a "kicker" line that gives the
    # title block a designed entry point.
    add_filled_rect(slide, left, 2.05, 0.9, 0.06, colour(brand, "accent"))

    # Title -- vertically centred in a generous block, with line spacing
    # tuned for long multi-word titles that wrap.
    add_text_box(
        slide, brand, left, 2.3, width, 1.9,
        as_text(spec.get("title"), "Untitled"),
        "header", sizes["title"], colour(brand, "ink"),
        align=align, bold=True, is_header=True,
    )
    # Promise reads at a clearly-secondary size, with comfortable wrap leading.
    add_multiline(
        slide, brand, left, 4.3, width, 1.4,
        [as_text(spec.get("promise"))],
        "body", sizes["heading"], colour(brand, "muted"),
        align=align, line_spacing=1.25,
    )
    add_text_box(
        slide, brand, left, 5.85, width, 0.5,
        as_text(spec.get("byline")),
        "body", sizes["body"], colour(brand, "accent"),
        align=align, bold=True,
    )

# === layouts/section_divider.py ===
"""statement family / section_divider variant: accent band, numeral, title."""

def _render_section_divider(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    accent_hex = accent_for_role(brand, spec.get("accent_role"))

    # Top accent band -- height scales with motif weight, colour from
    # accent_role so divider can signal section tone (e.g. critical section
    # opens with the warning accent rather than the deck's dominant accent).
    band_h = round(0.8 + 0.3 * brand["motif"]["weight"], 3)
    add_filled_rect(slide, 0, 0, geom["slide_w"], band_h, accent_hex)

    left = geom["body_left"]
    width = geom["body_w"]
    align = "center" if geom["alignment"] == "centered" else "left"

    # section_number coerced to int BEFORE any :02d-style format.
    number = coerce_int(spec.get("section_number"), 0)
    add_text_box(
        slide, brand, left, band_h + 0.6, 4.5, 2.5,
        "%02d" % number,
        "header", sizes["section_number"], accent_hex,
        bold=True, align=align,
    )
    # Thin rule between numeral and title for designed structure.
    rule_y = geom["slide_h"] - 2.9
    add_filled_rect(slide, left, rule_y, 1.4, 0.06, accent_hex)

    add_text_box(
        slide, brand, left, geom["slide_h"] - 2.6, width, 1.6,
        as_text(spec.get("title"), "Section"),
        "header", sizes["title"], colour(brand, "ink"),
        bold=True, is_header=True, align=align,
    )

# === layouts/single_statement.py ===
"""statement family / single_statement variant: one centred line, big type.

This is also the universal safe fallback when a variant is unknown.
"""



def _render_single_statement(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    accent_hex = accent_for_role(brand, spec.get("accent_role"))
    # Statement reads from `statement`, else any title-ish field.
    text = (
        spec.get("statement")
        or spec.get("title")
        or spec.get("heading")
        or spec.get("quote")
    )
    align = "left" if geom["alignment"] == "asymmetric" else "center"
    # Short accent rule above the statement -- visual entry point.
    rule_x = (geom["slide_w"] - 1.4) / 2 if align == "center" else geom["body_left"]
    add_filled_rect(slide, rule_x, 2.1, 1.4, 0.06, accent_hex)

    add_text_box(
        slide, brand, geom["body_left"], 2.5, geom["body_w"], 2.7,
        as_text(text, ""),
        "header", sizes["statement"], colour(brand, "ink"),
        align=align, anchor="middle",
    )

# === layouts/recap.py ===
"""statement family / recap variant: bulleted takeaways + next-step pointer."""



def _render_recap(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    accent_hex = accent_for_role(brand, spec.get("accent_role"))

    draw_heading(slide, brand, geom, "Takeaways",
                 accent_role=spec.get("accent_role", "neutral"))

    content_top = geom["content_top"] + 0.15
    left = geom["body_left"]
    width = geom["body_w"]
    split = left + width * 0.62
    height = geom["content_bottom"] - content_top

    takeaways = spec.get("takeaways")
    if not isinstance(takeaways, list):
        takeaways = [as_text(takeaways)] if takeaways else []
    # Accent-coloured bullets on takeaway list -- subtle finish that ties
    # the slide back to the deck identity.
    add_multiline(
        slide, brand, left, content_top, width * 0.58, height,
        [as_text(t) for t in takeaways],
        "body", sizes["body"] + 1, colour(brand, "ink"),
        bullet=True, bullet_colour_hex=accent_hex,
        line_spacing=1.35,
    )

    # Next-step block sits on a surface card with a top accent rule.
    card_x, card_y = split, content_top
    card_w, card_h = width * 0.38, height
    add_filled_rect(slide, card_x, card_y, card_w, card_h, colour(brand, "surface"))
    add_filled_rect(slide, card_x, card_y, card_w, 0.08, accent_hex)
    pad = 0.3
    add_text_box(
        slide, brand, card_x + pad, card_y + 0.3, card_w - 2 * pad, 0.5,
        "Next step", "header", sizes["pillar_heading"],
        accent_hex, bold=True, is_header=True,
    )
    add_text_box(
        slide, brand, card_x + pad, card_y + 0.95, card_w - 2 * pad,
        card_h - 1.1,
        as_text(spec.get("next_step")),
        "body", sizes["body"] + 1, colour(brand, "ink"),
    )

# === layouts/stat_callout.py ===
"""focus family / stat_callout variant: hero number with unit + context."""


# Units that read naturally INLINE with the number (no line break).
# Words like "hours", "kg", "students" still stack on their own line.
_INLINE_UNITS = {"%", "°", "×", "x", "+", "-", "K", "M", "B", "k", "m"}


def _render_stat_callout(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    left = geom["body_left"]
    width = geom["body_w"]
    align = "left" if geom["alignment"] == "asymmetric" else "center"
    accent_hex = accent_for_role(brand, spec.get("accent_role"))

    stat = as_text(spec.get("stat"))
    unit = as_text(spec.get("unit"))

    # If unit is a single symbol (%, °, +, ×, K, M, B), render inline with the
    # number. Otherwise stack it underneath so e.g. "12 students" stays clear.
    inline_unit = unit and (unit.strip() in _INLINE_UNITS or len(unit.strip()) == 1)
    if inline_unit:
        # Stat and inline unit as one display block. Unit gets ~55% of the
        # stat size so it reads as a partner glyph, not a sibling number.
        # We approximate this with two text boxes side-by-side, anchored by
        # the unit length.
        # Simple path: render "stat + unit" together at the hero size, then
        # null the unit so the stacked-unit branch is skipped.
        hero_text = stat + unit
        add_text_box(
            slide, brand, left, 2.0, width, 2.6,
            hero_text,
            "header", sizes["stat"], accent_hex,
            align=align, anchor="middle", bold=True,
        )
        context_top = 4.5
    else:
        add_text_box(
            slide, brand, left, 1.5, width, 2.4,
            stat,
            "header", sizes["stat"], accent_hex,
            align=align, anchor="middle", bold=True,
        )
        if unit:
            add_text_box(
                slide, brand, left, 3.85, width, 0.8,
                unit,
                "header", sizes["heading"], colour(brand, "ink"),
                align=align, anchor="top", bold=True,
            )
            context_top = 4.8
        else:
            context_top = 4.1

    add_text_box(
        slide, brand, left, context_top, width, 1.8,
        as_text(spec.get("context")),
        "body", sizes["heading"], colour(brand, "muted"),
        align=align, anchor="top",
    )

# === layouts/quote.py ===
"""focus family / quote variant: oversized quote mark + body + attribution."""



def _render_quote(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    left = geom["body_left"]
    width = geom["body_w"]
    accent_hex = accent_for_role(brand, spec.get("accent_role"))

    # Oversized opening mark in the accent colour.
    add_text_box(
        slide, brand, left, 0.6, 2.2, 1.7,
        "\u201C",
        "header", int(sizes["section_number"] * 1.6), accent_hex,
        bold=True,
    )
    # Quote body wrapped in proper curly quotes (open already drawn oversized,
    # but include both inline too so the body reads as a complete sentence).
    body_quote = as_text(spec.get("quote"))
    if body_quote and not body_quote.endswith("\u201D"):
        body_quote = body_quote.rstrip(".") + "\u201D"
    add_text_box(
        slide, brand, left, 2.1, width, 3.2,
        body_quote,
        "header", sizes["heading"], colour(brand, "ink"),
    )
    # Em-dash before attribution (designed convention, not "- ").
    add_text_box(
        slide, brand, left, 5.5, width, 0.7,
        "\u2014 " + as_text(spec.get("attribution")),
        "body", sizes["body"], colour(brand, "muted"),
    )

# === layouts/two_up.py ===
"""split family / two_up variant: heading, body left, configurable right block.

right_block_type is one of quote | stat | diagram (chart treated as diagram).
"""



def _heading(slide, spec, brand, geom):
    """Kept for backward-compat; new code calls draw_heading directly."""
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))


def _right_quote(slide, box, content, brand, accent_hex):
    sizes = brand["sizes_pt"]
    left, top, width, height = box
    add_text_box(
        slide, brand, left, top, width, height - 0.8,
        "\u201C" + as_text(content.get("quote")) + "\u201D",
        "header", sizes["pillar_heading"] + 4, colour(brand, "ink"),
    )
    add_text_box(
        slide, brand, left, top + height - 0.7, width, 0.5,
        "\u2014 " + as_text(content.get("attribution")),
        "body", sizes["body"], colour(brand, "muted"),
    )


def _right_stat(slide, box, content, brand, accent_hex):
    sizes = brand["sizes_pt"]
    left, top, width, height = box
    stat = as_text(content.get("stat"))
    unit = as_text(content.get("unit"))

    # Same inline-symbol rule as _render_stat_callout.
    inline_unit = unit and (unit.strip() in _INLINE_UNITS or len(unit.strip()) == 1)
    if inline_unit:
        add_text_box(
            slide, brand, left, top, width, 1.7,
            stat + unit,
            "header", sizes["stat"] - 24, accent_hex,
            bold=True, anchor="middle", align="center",
        )
        ctx_top = top + 1.8
    else:
        add_text_box(
            slide, brand, left, top, width, 1.5,
            stat,
            "header", sizes["stat"] - 24, accent_hex,
            bold=True, anchor="middle", align="center",
        )
        if unit:
            add_text_box(
                slide, brand, left, top + 1.5, width, 0.5,
                unit,
                "header", sizes["body"] + 2, colour(brand, "ink"),
                bold=True, anchor="top", align="center",
            )
        ctx_top = top + 2.2

    add_text_box(
        slide, brand, left, ctx_top, width, height - (ctx_top - top),
        as_text(content.get("context")),
        "body", sizes["body"], colour(brand, "muted"),
        align="center",
    )


def _right_diagram(slide, box, content, brand, accent_hex):
    left, top, width, height = box
    add_text_box(
        slide, brand, left, top, width, height,
        as_text(content.get("caption"), "[diagram]"),
        "body", brand["sizes_pt"]["body"], colour(brand, "muted"),
        align="center", anchor="middle",
    )


_RIGHT = {
    "quote": _right_quote,
    "stat": _right_stat,
    "diagram": _right_diagram,
    "chart": _right_diagram,
}


def _render_two_up(slide, spec, brand, geom):
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))
    sizes = brand["sizes_pt"]
    accent_hex = accent_for_role(brand, spec.get("accent_role"))
    top = geom["content_top"] + 0.15
    height = geom["content_bottom"] - top
    gap = geom["gap"]
    left = geom["body_left"]
    col_w = (geom["body_w"] - gap) / 2.0

    add_text_box(
        slide, brand, left, top, col_w, height,
        as_text(spec.get("body")),
        "body", sizes["body"], colour(brand, "ink"),
    )

    right_left = left + col_w + gap
    # Surface-tinted card with a top accent rule -- ties right block to deck.
    add_filled_rect(slide, right_left, top, col_w, height, colour(brand, "surface"))
    add_filled_rect(slide, right_left, top, col_w, 0.08, accent_hex)

    rtype = spec.get("right_block_type", "diagram")
    content = spec.get("right_block_content")
    if not isinstance(content, dict):
        content = {}
    pad = 0.3
    box = (right_left + pad, top + 0.4, col_w - 2 * pad, height - 0.5)
    _RIGHT.get(rtype, _right_diagram)(slide, box, content, brand, accent_hex)

# === layouts/compare_contrast.py ===
"""split family / compare_contrast variant: heading + two labelled columns."""



def _render_compare_contrast(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    # Heading carries the deck's neutral accent rule.
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))

    top = geom["content_top"] + 0.15
    height = geom["content_bottom"] - top
    gap = geom["gap"]
    left = geom["body_left"]
    col_w = (geom["body_w"] - gap) / 2.0

    # Per-side accent_role drives the column accent. Writers can flag:
    #   left_accent_role:  "warning" / "critical"  (e.g. the "Bystander" side)
    #   right_accent_role: "positive"              (e.g. the "Upstander" side)
    # Falling back to (left=neutral, right=positive) is safer than the old
    # hard-coded (left=accent, right=accent_2) which painted *both* sides
    # with deck accents regardless of meaning.
    side_roles = {
        "left": spec.get("left_accent_role", "neutral"),
        "right": spec.get("right_accent_role", "positive"),
    }

    for idx, side in enumerate(("left", "right")):
        x = left + idx * (col_w + gap)
        col_accent = accent_for_role(brand, side_roles[side])

        # 0.08" accent bar across the top of the column. Strong visual signature
        # that the two columns are semantically different.
        add_filled_rect(slide, x, top, col_w, 0.08, col_accent)

        # Label in the column's accent.
        add_text_box(
            slide, brand, x, top + 0.18, col_w, 0.6,
            as_text(spec.get(side + "_label")),
            "header", sizes["pillar_heading"], col_accent,
            bold=True, is_header=True,
        )
        points = spec.get(side + "_points")
        if not isinstance(points, list):
            points = [as_text(points)] if points else []
        # Bullets sit in the column's accent colour -- subtle but designerly.
        add_multiline(
            slide, brand, x, top + 0.95, col_w, height - 1.0,
            [as_text(p) for p in points],
            "body", sizes["body"], colour(brand, "ink"),
            bullet=True, bullet_colour_hex=col_accent,
            line_spacing=1.25,
        )

# === layouts/pillars.py ===
"""grid family / pillars_3 + pillars_4 variants: heading + N equal columns.

Count is enforced by VARIANT_ELEMENT_COUNTS: pillars_3 -> exactly 3,
pillars_4 -> exactly 4. A wrong count is padded/truncated, never raised.
"""

def _empty_pillar():
    return {"icon_glyph": "", "heading": "", "body": "", "accent_role": "neutral"}


def _render_pillars(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    variant = spec.get("variant", "pillars_3")
    count = 4 if variant == "pillars_4" else 3

    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))

    pillars = coerce_count(spec.get("pillars"), count, count, _empty_pillar)

    top = geom["content_top"] + 0.15
    height = geom["content_bottom"] - top
    gap = max(geom["gap"], 0.25)
    left = geom["body_left"]
    col_w = (geom["body_w"] - gap * (count - 1)) / count
    body_size = sizes["body"] if count == 3 else max(11, sizes["body"] - 3)

    # Cards fill the full body height -- no more dead space underneath. The
    # top accent stripe is the per-pillar signature; we drop the old icon-glyph
    # square because LLM writers consistently fell back to meaningless single
    # letters. A per-pillar accent_role lets writers paint specific pillars
    # warning/positive when the meaning warrants.
    for i, pillar in enumerate(pillars):
        if not isinstance(pillar, dict):
            pillar = _empty_pillar()
        x = left + i * (col_w + gap)

        pillar_accent = accent_for_role(brand, pillar.get("accent_role"))

        # Card body (surface fill).
        add_filled_rect(slide, x, top, col_w, height, colour(brand, "surface"))
        # Top accent stripe -- the pillar's identity.
        add_filled_rect(slide, x, top, col_w, 0.10, pillar_accent)

        # Heading.
        pad_x = 0.25
        head_top = top + 0.35
        add_text_box(
            slide, brand, x + pad_x, head_top, col_w - 2 * pad_x, 0.7,
            as_text(pillar.get("heading")),
            "header", sizes["pillar_heading"], colour(brand, "ink"),
            bold=True, is_header=True,
        )
        # Body.
        body_top = head_top + 0.85
        add_multiline(
            slide, brand, x + pad_x, body_top, col_w - 2 * pad_x,
            top + height - body_top - 0.25,
            [as_text(pillar.get("body"))],
            "body", body_size, colour(brand, "ink"),
            line_spacing=1.25,
        )

# === layouts/quadrant_grid.py ===
"""grid family / quadrant_grid variant: heading + 2x2 grid of numbered cells.

VARIANT_ELEMENT_COUNTS: exactly 4 cells -- padded/truncated, never raised.
"""

def _empty_cell():
    return {"number": "", "label": "", "body": "", "accent_role": "neutral"}


def _render_quadrant_grid(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))

    cells = coerce_count(spec.get("cells"), 4, 4, _empty_cell)

    top = geom["content_top"] + 0.15
    height = geom["content_bottom"] - top
    gap = max(geom["gap"], 0.25)
    left = geom["body_left"]
    cw = (geom["body_w"] - gap) / 2.0
    ch = (height - gap) / 2.0

    positions = [
        (left, top), (left + cw + gap, top),
        (left, top + ch + gap), (left + cw + gap, top + ch + gap),
    ]

    for (x, y), cell in zip(positions, cells):
        if not isinstance(cell, dict):
            cell = _empty_cell()
        cell_accent = accent_for_role(brand, cell.get("accent_role"))

        add_rounded_rect(slide, x, y, cw, ch, colour(brand, "surface"),
                         line_hex=cell_accent)
        add_text_box(
            slide, brand, x + 0.2, y + 0.12, 1.1, 1.0,
            as_text(cell.get("number")),
            "header", max(20, sizes["heading"] + 8), cell_accent,
            bold=True,
        )
        add_text_box(
            slide, brand, x + 1.3, y + 0.25, cw - 1.5, 0.6,
            as_text(cell.get("label")),
            "header", sizes["pillar_heading"], colour(brand, "ink"),
            bold=True, is_header=True,
        )
        add_multiline(
            slide, brand, x + 0.3, y + 1.0, cw - 0.6, ch - 1.2,
            [as_text(cell.get("body"))],
            "body", sizes["body"], colour(brand, "ink"),
            line_spacing=1.25,
        )

# === layouts/process_flow.py ===
"""flow family / process_flow variant: heading + 4-6 step boxes in a row.

VARIANT_ELEMENT_COUNTS: 4-6 steps. Fewer/more -> padded/truncated, never raised.
"""

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

def _empty_step():
    return {"number": "", "label": "", "detail": "", "accent_role": "neutral"}


def _render_process_flow(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))

    steps = coerce_count(spec.get("steps"), 4, 6, _empty_step)
    n = len(steps)

    top = geom["content_top"] + 0.2
    gap = max(geom["gap"], 0.2)
    left = geom["body_left"]
    # Chevron tail width -- reserve from the right edge of each box.
    chev_w = 0.22
    box_w = (geom["body_w"] - gap * (n - 1)) / n
    box_h = 1.9
    deck_accent = accent_for_role(brand, spec.get("accent_role"))

    for i, step in enumerate(steps):
        if not isinstance(step, dict):
            step = _empty_step()
        x = left + i * (box_w + gap)

        # Per-step accent_role lets writers escalate one step (e.g. the
        # "danger" or "success" terminus). Default uses the deck accent.
        step_accent = accent_for_role(brand, step.get("accent_role"))
        on_accent = readable_on(brand, step_accent)

        # Chevron-shaped box for all but the last; rounded rect for the last.
        is_last = (i == n - 1)
        if is_last:
            rect = add_rounded_rect(slide, x, top, box_w, box_h, step_accent)
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.PENTAGON,  # right-pointing chevron-like shape
                Inches(x), Inches(top), Inches(box_w), Inches(box_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(step_accent)
            shape.line.fill.background()
            shape.shadow.inherit = False
            rect = shape

        # Step number badge in the top-left corner of each box.
        num = coerce_int(step.get("number"), i + 1)
        add_text_box(
            slide, brand, x + 0.15, top + 0.12, 0.9, 0.5,
            "%02d" % num,
            "header", max(14, sizes["body"] + 4), on_accent,
            bold=True,
        )

        # Step label fills the rest of the box, vertically centred.
        tf = rect.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "\n\n" + apply_header_case(brand, as_text(step.get("label")))
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            run = p.runs[0]
            run.font.name = brand["fonts"]["header"]
            run.font.size = Pt(sizes["pillar_heading"])
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(on_accent)

        # Detail text below the step box.
        add_text_box(
            slide, brand, x, top + box_h + 0.2, box_w,
            geom["content_bottom"] - top - box_h - 0.2,
            as_text(step.get("detail")),
            "body", max(11, sizes["body"] - 3), colour(brand, "muted"),
            align="center",
        )

# === layouts/chart.py ===
"""data family / chart variant: heading + native PPTX chart + caption."""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches



_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
}


def _render_chart(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))

    top = geom["content_top"] + 0.15
    height = geom["content_bottom"] - top - 0.5
    chart_type = _CHART_TYPE_MAP.get(
        spec.get("chart_type", "bar"), XL_CHART_TYPE.BAR_CLUSTERED
    )

    data = spec.get("chart_data")
    if not isinstance(data, dict):
        data = {}
    categories = data.get("categories")
    if not isinstance(categories, list) or not categories:
        categories = ["A", "B"]
    series = data.get("series")
    if not isinstance(series, list) or not series:
        series = [{"name": "Series 1", "values": [1] * len(categories)}]

    chart_data = CategoryChartData()
    chart_data.categories = [as_text(c) for c in categories]
    for s in series:
        if not isinstance(s, dict):
            continue
        values = s.get("values")
        if not isinstance(values, list):
            values = [0] * len(categories)
        # Pad/truncate values to category count.
        values = (values + [0] * len(categories))[:len(categories)]
        nums = []
        for v in values:
            try:
                nums.append(float(v))
            except (TypeError, ValueError):
                nums.append(0.0)
        chart_data.add_series(as_text(s.get("name"), "Series"), nums)

    slide.shapes.add_chart(
        chart_type,
        Inches(geom["body_left"]), Inches(top),
        Inches(geom["body_w"]), Inches(height),
        chart_data,
    )
    add_text_box(
        slide, brand, geom["body_left"], top + height + 0.05,
        geom["body_w"], 0.5,
        as_text(spec.get("caption")),
        "body", sizes["caption"], colour(brand, "muted"),
    )

# === layouts/data_table.py ===
"""data family / data_table variant: heading + native PPTX table.

Limits: <=5 rows x <=4 cols. Header row tinted with accent; emphasis cells
get a surface fill. Over-limit input is truncated, never raised.
"""

from pptx.util import Inches, Pt

_MAX_ROWS = 5
_MAX_COLS = 4


def _render_data_table(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))
    accent = accent_for_role(brand, spec.get("accent_role"))

    columns = spec.get("columns")
    if not isinstance(columns, list) or not columns:
        columns = ["Column"]
    columns = [as_text(c) for c in columns[:_MAX_COLS]]
    n_cols = len(columns)

    rows = spec.get("rows")
    if not isinstance(rows, list):
        rows = []
    rows = rows[:_MAX_ROWS]
    norm_rows = []
    for row in rows:
        if not isinstance(row, (list, tuple)):
            row = [as_text(row)]
        cells = [as_text(v) for v in row][:n_cols]
        cells += [""] * (n_cols - len(cells))
        norm_rows.append(cells)

    top = geom["content_top"] + 0.15
    table = slide.shapes.add_table(
        rows=len(norm_rows) + 1, cols=n_cols,
        left=Inches(geom["body_left"]), top=Inches(top),
        width=Inches(geom["body_w"]),
        height=Inches(geom["content_bottom"] - top),
    ).table

    on_accent = readable_on(brand, accent)
    ink = colour(brand, "ink")
    surface = colour(brand, "surface")

    for c, col in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(accent)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = brand["fonts"]["header"]
                run.font.size = Pt(sizes["caption"] + 3)
                run.font.bold = True
                run.font.color.rgb = hex_to_rgb(on_accent)

    emphasis = set()
    for pair in spec.get("emphasis_cells", []) or []:
        if isinstance(pair, (list, tuple)) and len(pair) == 2:
            emphasis.add((pair[0], pair[1]))

    # Zebra striping: every other row gets a surface tint. Emphasis cells
    # override with the explicit emphasis fill.
    for r, row in enumerate(norm_rows):
        is_zebra = (r % 2 == 1)
        for c, value in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = value
            if (r, c) in emphasis:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(accent)
                cell_text_hex = on_accent
            elif is_zebra:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(surface)
                cell_text_hex = ink
            else:
                cell_text_hex = ink
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = brand["fonts"]["body"]
                    run.font.size = Pt(sizes["caption"] + 1)
                    run.font.color.rgb = hex_to_rgb(cell_text_hex)

# === layouts/stat_doughnut.py ===
"""data family / stat_doughnut variant: heading + 3 side-by-side doughnuts.

VARIANT_ELEMENT_COUNTS: exactly 3 doughnuts -- padded/truncated, never raised.
"""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches



def _empty_doughnut():
    return {"label": "", "percent": 0, "accent_role": "neutral"}


def _pct(value):
    try:
        n = float(value)
    except (TypeError, ValueError):
        return 0.0
    return max(0.0, min(100.0, n))


def _render_stat_doughnut(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    draw_heading(slide, brand, geom, spec.get("heading"),
                 accent_role=spec.get("accent_role", "neutral"))

    doughnuts = coerce_count(spec.get("doughnuts"), 3, 3, _empty_doughnut)
    n = 3
    top = geom["content_top"] + 0.15
    left = geom["body_left"]
    size = min(3.2, (geom["body_w"] - 0.6) / n)
    gap = (geom["body_w"] - size * n) / (n - 1) if n > 1 else 0

    # accent_role -> palette role (warning/critical map to accent_2).
    role_map = {
        "positive": "accent", "neutral": "accent",
        "warning": "accent_2", "critical": "accent_2",
    }

    for i, d in enumerate(doughnuts):
        if not isinstance(d, dict):
            d = _empty_doughnut()
        x = left + i * (size + gap)
        pct = _pct(d.get("percent"))

        chart_data = CategoryChartData()
        chart_data.categories = [as_text(d.get("label"), "share"), "rest"]
        chart_data.add_series("share", [pct, 100 - pct])
        slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(x), Inches(top), Inches(size), Inches(size),
            chart_data,
        )
        role = role_map.get(d.get("accent_role", "neutral"), "accent")
        add_text_box(
            slide, brand, x, top + size + 0.05, size, 0.55,
            "%d%%" % int(round(pct)),
            "header", sizes["pillar_heading"] + 6, colour(brand, role),
            bold=True, align="center",
        )
        add_text_box(
            slide, brand, x, top + size + 0.65, size, 0.7,
            as_text(d.get("label")),
            "body", sizes["body"], colour(brand, "ink"),
            align="center",
        )

# === layouts/callout_box.py ===
"""callout family / callout_box variant: severity-keyed tinted container.

severity is one of info | positive | warning | critical. It maps onto the
generative palette (not a fixed brand palette): info -> muted, positive ->
accent, warning/critical -> accent_2.
"""

_SEVERITY_ROLE = {
    "info": "muted",
    "positive": "accent",
    "warning": "accent_2",
    "critical": "accent_2",
}


# Severity glyphs read at scale -- single Unicode characters that ship with
# every font we use. Keeps the badge visible even when the font lacks
# emoji support.
_SEVERITY_GLYPH = {
    "info":     "i",
    "positive": "+",
    "warning":  "!",
    "critical": "!",
}


def _render_callout_box(slide, spec, brand, geom):
    sizes = brand["sizes_pt"]
    severity = spec.get("severity", "info")
    fill = colour(brand, _SEVERITY_ROLE.get(severity, "muted"))
    text_hex = readable_on(brand, fill)

    # Container inset from the body region.
    inset = 0.4
    left = geom["body_left"] + inset
    top = geom["content_top"]
    width = geom["body_w"] - 2 * inset
    height = geom["content_bottom"] - top - 0.3

    add_rounded_rect(slide, left, top, width, height, fill)

    # Severity badge -- a small circle in the top-left with a single glyph.
    # Reads as "this is a warning/info" without needing to load an icon font.
    badge_size = 0.65
    badge_x = left + 0.4
    badge_y = top + 0.4
    badge_fill = text_hex  # use the readable text colour for the badge fill
    add_oval(slide, badge_x, badge_y, badge_size, badge_size, badge_fill)
    add_text_box(
        slide, brand, badge_x, badge_y - 0.02, badge_size, badge_size,
        _SEVERITY_GLYPH.get(severity, "i"),
        "header", int(sizes["heading"] * 0.9), fill,
        bold=True, align="center", anchor="middle",
    )

    pad = 0.45
    head_left = badge_x + badge_size + 0.3
    add_text_box(
        slide, brand, head_left, top + pad,
        width - (head_left - left) - pad, 0.9,
        as_text(spec.get("heading")),
        "header", sizes["heading"], text_hex,
        bold=True, is_header=True,
    )
    add_multiline(
        slide, brand, left + pad, top + pad + 1.2,
        width - 2 * pad, height - pad - 1.5,
        [as_text(spec.get("body"))],
        "body", sizes["body"] + 2, text_hex,
        line_spacing=1.4,
    )

# === renderer.py ===
"""render_deck -- token-driven entry point.

Signature: render_deck(deck_spec, design_tokens) -> pptx.Presentation.

  * `deck_spec`  -- {deck_title, footer_text, narrative_arc, sections, slides}.
    Each slide carries layout_family, variant, position, branch_index,
    accent_role, presenter_notes, optional per-slide `composition` override,
    plus the variant content fields (VARIANT_CONTENT in schema.py).
  * `design_tokens` -- the generative visual system (DESIGN_TOKEN_SCHEMA).

Dispatch is on `variant`. An unknown variant renders the safe
`single_statement` fallback rather than crashing. The renderer returns the
Presentation object; the caller saves it.
"""

from pptx import Presentation
from pptx.util import Inches




# variant -> layout render function. 16 variants across 7 families.
_RENDERERS = {
    "title": _render_title,
    "section_divider": _render_section_divider,
    "single_statement": _render_single_statement,
    "recap": _render_recap,
    "stat_callout": _render_stat_callout,
    "quote": _render_quote,
    "two_up": _render_two_up,
    "compare_contrast": _render_compare_contrast,
    "pillars_3": _render_pillars,
    "pillars_4": _render_pillars,
    "quadrant_grid": _render_quadrant_grid,
    "process_flow": _render_process_flow,
    "chart": _render_chart,
    "data_table": _render_data_table,
    "stat_doughnut": _render_stat_doughnut,
    "callout_box": _render_callout_box,
}

# Variants that should NOT receive a footer (full-bleed display slides).
_NO_FOOTER = {"title", "section_divider"}
# Variants that should NOT receive the motif (full-bleed display slides).
_NO_MOTIF = {"title", "section_divider"}


def _slide_brand(design_tokens, base_brand, slide_spec):
    """Resolve the brand for a slide, applying any per-slide composition override.

    An override carries raw token values (e.g. margin="wide"), so it is merged
    at the token level and re-resolved -- never injected into a resolved brand.
    """
    override = slide_spec.get("composition")
    if not isinstance(override, dict) or not override:
        return base_brand
    tokens = design_tokens if isinstance(design_tokens, dict) else {}
    merged = dict(tokens)
    comp = dict(tokens.get("composition", {}) or {})
    for key in ("alignment", "density", "margin", "whitespace"):
        if key in override:
            comp[key] = override[key]
    merged["composition"] = comp
    return resolve_brand(merged)


def _draw_footer(slide, brand, geom, footer_text, position):
    """Two-part footer: footer_text in muted, page number in accent.
    Keeps the deck identity visible without dominating any slide.
    """
    if not footer_text:
        return
    n = coerce_int(position, 0)
    label = as_text(footer_text)
    sizes = brand["sizes_pt"]
    y = geom["slide_h"] - 0.45

    # Footer label (left of separator) in muted.
    add_text_box(
        slide, brand, geom["margin"], y,
        geom["content_w"] - 0.4, 0.35,
        label, "body", sizes["caption"],
        colour(brand, "muted"), align="right",
    )
    # Page number in accent at the far right -- small designed touch that
    # makes the deck feel produced rather than auto-generated.
    if n:
        add_text_box(
            slide, brand, geom["slide_w"] - geom["margin"] - 0.6, y,
            0.55, 0.35,
            "%02d" % n, "body", sizes["caption"],
            colour(brand, "accent"), align="right", bold=True,
        )


def render_deck(deck_spec, design_tokens):
    """Render a deck from a deck_spec + design_tokens. Returns a Presentation."""
    deck_spec = deck_spec if isinstance(deck_spec, dict) else {}
    brand = resolve_brand(design_tokens)

    prs = Presentation()
    dims = brand["slide_dimensions_in"]
    prs.slide_width = Inches(dims["width"])
    prs.slide_height = Inches(dims["height"])
    blank = prs.slide_layouts[6]

    slides = deck_spec.get("slides")
    if not isinstance(slides, list):
        slides = []
    slides = sorted(slides, key=lambda s: coerce_int(s.get("position"), 0))
    footer_text = as_text(deck_spec.get("footer_text"))

    for slide_spec in slides:
        if not isinstance(slide_spec, dict):
            continue
        slide = prs.slides.add_slide(blank)
        slide_brand = _slide_brand(design_tokens, brand, slide_spec)
        geom = compute_geometry(slide_brand)

        fill_background(slide, slide_brand)

        variant = slide_spec.get("variant", "single_statement")
        if variant not in _NO_MOTIF:
            draw_motif(slide, slide_brand)

        renderer = _RENDERERS.get(variant, _render_single_statement)
        try:
            renderer(slide, slide_spec, slide_brand, geom)
        except Exception:
            # Last-ditch guard: never crash a whole deck on one bad slide.
            _render_single_statement(slide, slide_spec, slide_brand, geom)

        if variant not in _NO_FOOTER:
            _draw_footer(slide, slide_brand, geom, footer_text,
                         slide_spec.get("position"))

        notes = as_text(slide_spec.get("presenter_notes"))
        slide.notes_slide.notes_text_frame.text = notes

    return prs
